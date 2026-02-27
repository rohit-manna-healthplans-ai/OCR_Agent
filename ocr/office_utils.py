"""
Support for DOC, DOCX, PPT, PPTX in universal OCR.

- DOCX/PPTX: parsed natively (python-docx, python-pptx); extract text and images, OCR images.
- DOC/PPT: converted to PDF via LibreOffice (soffice --headless), then processed as PDF.
"""

from __future__ import annotations

import io
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import List, Optional, Tuple

import numpy as np
from PIL import Image


def docx_extract_text(path: str) -> str:
    """Extract all paragraph text from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        return ""
    doc = Document(path)
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip()).strip()


def docx_extract_images(path: str, max_images: int = 50) -> List[np.ndarray]:
    """Extract embedded images from DOCX (word/media/*) as RGB numpy arrays. Skips WMF/EMF if PIL can't open."""
    out: List[np.ndarray] = []
    try:
        with zipfile.ZipFile(path, "r") as z:
            for info in z.filelist:
                if not info.filename.startswith("word/media/"):
                    continue
                if len(out) >= max_images:
                    break
                ext = (Path(info.filename).suffix or "").lower()
                if ext in (".wmf", ".emf"):
                    continue
                try:
                    data = z.read(info.filename)
                    img = Image.open(io.BytesIO(data)).convert("RGB")
                    out.append(np.array(img))
                except Exception:
                    continue
    except Exception:
        pass
    return out


def docx_load(path: str, max_pages: int = 50) -> Tuple[str, List[np.ndarray]]:
    """Return (full text, list of image arrays) for DOCX. Images limited by max_pages."""
    text = docx_extract_text(path)
    images = docx_extract_images(path, max_images=max_pages)
    return text, images


def pptx_extract_slides(path: str, max_slides: int = 50) -> List[Tuple[str, List[np.ndarray]]]:
    """For each slide, return (slide text, list of image arrays)."""
    try:
        from pptx import Presentation
        from pptx.shapes.picture import Picture
    except ImportError:
        return []

    out: List[Tuple[str, List[np.ndarray]]] = []
    prs = Presentation(path)
    for slide in prs.slides[:max_slides]:
        texts: List[str] = []
        images: List[np.ndarray] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
            if isinstance(shape, Picture) and hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob)).convert("RGB")
                    images.append(np.array(img))
                except Exception:
                    pass
        text = "\n".join(t for t in texts if t).strip()
        out.append((text, images))
    return out


def convert_office_to_pdf(input_path: str) -> Optional[str]:
    """
    Convert .doc or .ppt to PDF using LibreOffice (soffice --headless).
    Returns path to temporary PDF file, or None if conversion fails.
    """
    p = Path(input_path)
    suffix = p.suffix.lower()
    if suffix not in (".doc", ".ppt"):
        return None

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        pdf_path = tmp.name
    out_dir = Path(pdf_path).parent
    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(out_dir),
                str(input_path),
            ],
            capture_output=True,
            timeout=120,
            check=True,
        )
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        try:
            Path(pdf_path).unlink(missing_ok=True)
        except Exception:
            pass
        return None

    # LibreOffice writes file with same stem, e.g. doc.pdf
    converted = out_dir / (p.stem + ".pdf")
    if converted.exists():
        return str(converted)
    return None
